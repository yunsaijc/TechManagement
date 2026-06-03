from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path("/home/tdkx/ljh/Tech/系统平台介绍@0412 - 副本.pptx")
OUT = Path("/home/tdkx/ljh/Tech/系统平台介绍@0412 - 优化版.pptx")

NAVY = RGBColor(15, 32, 58)
TEAL = RGBColor(16, 111, 112)
BLUE = RGBColor(34, 87, 190)
GOLD = RGBColor(196, 149, 63)
SLATE = RGBColor(85, 98, 118)
LIGHT_BG = RGBColor(244, 247, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(217, 225, 236)

SECTION_STYLES = {
    "网络舆情与内容安全治理平台": ("Platform 01", TEAL),
    "内容安全认知博弈对抗平台": ("Platform 02", BLUE),
    "隐蔽网络数据获取与分析平台": ("Platform 03", GOLD),
}


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def hide_shape(shape) -> None:
    try:
        if hasattr(shape, "fill"):
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_BG
            shape.line.fill.background()
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()
    except Exception:
        return


def normalize_text(shape, *, font_size: int, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_header_band(slide, label: str, accent: RGBColor, slide_width: int) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, Inches(0.22))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.18), Inches(1.45), Inches(0.35))
    chip.fill.solid()
    chip.fill.fore_color.rgb = accent
    chip.line.fill.background()
    chip.text_frame.text = label
    normalize_text(chip, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_footer(slide, idx: int, total: int, accent: RGBColor) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(7.02), Inches(12.0), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(11.6), Inches(6.82), Inches(0.65), Inches(0.22))
    box.text_frame.text = f"{idx:02d}"
    normalize_text(box, font_size=10, color=accent, bold=True, align=PP_ALIGN.RIGHT)


def get_title_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = " ".join((shape.text or "").split())
        if not text or len(text) > 40:
            continue
        if shape.top > Inches(1.2):
            continue
        candidates.append(shape)
    if not candidates:
        return None
    return sorted(candidates, key=lambda s: (s.top, s.left))[0]


def style_title_shape(shape, accent: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.8)
    normalize_text(shape, font_size=22, color=NAVY, bold=True)


def style_body_text(shape) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text = " ".join((shape.text or "").split())
    if not text:
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    normalize_text(shape, font_size=16, color=SLATE)


def style_process_box(shape, accent: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.6)
    normalize_text(shape, font_size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def style_cover(slide) -> None:
    set_slide_bg(slide, NAVY)
    for shape in list(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = " ".join((shape.text or "").split())
            if text == "网络空间舆情治理与内容安全博弈对抗":
                shape.left = Inches(0.9)
                shape.top = Inches(2.2)
                shape.width = Inches(8.8)
                shape.height = Inches(1.1)
                shape.fill.solid()
                shape.fill.fore_color.rgb = NAVY
                shape.line.fill.background()
                normalize_text(shape, font_size=26, color=WHITE, bold=True)
            elif text:
                shape.fill.solid()
                shape.fill.fore_color.rgb = NAVY
                shape.line.fill.background()
                normalize_text(shape, font_size=12, color=RGBColor(174, 189, 214), bold=False)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(1.7), Inches(1.45), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    sub = slide.shapes.add_textbox(Inches(0.95), Inches(3.55), Inches(4.2), Inches(0.45))
    sub.text_frame.text = "系统平台介绍汇报"
    normalize_text(sub, font_size=16, color=RGBColor(198, 208, 226), bold=False)


def style_agenda(slide) -> None:
    set_slide_bg(slide, LIGHT_BG)
    title = get_title_shape(slide)
    if title:
        title.left = Inches(0.75)
        title.top = Inches(0.6)
        title.width = Inches(2.6)
        title.height = Inches(0.5)
        title.fill.solid()
        title.fill.fore_color.rgb = LIGHT_BG
        title.line.fill.background()
        normalize_text(title, font_size=24, color=NAVY, bold=True)
    add_header_band(slide, "Overview", TEAL, 12192000)
    add_footer(slide, 2, 31, TEAL)


def style_end(slide) -> None:
    set_slide_bg(slide, NAVY)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = " ".join((shape.text or "").split())
            if "感谢聆听" in text:
                shape.fill.solid()
                shape.fill.fore_color.rgb = NAVY
                shape.line.fill.background()
                shape.left = Inches(3.25)
                shape.top = Inches(2.8)
                shape.width = Inches(6.0)
                shape.height = Inches(0.9)
                normalize_text(shape, font_size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            else:
                hide_shape(shape)
    sub = slide.shapes.add_textbox(Inches(4.0), Inches(3.9), Inches(4.7), Inches(0.35))
    sub.text_frame.text = "期待交流指正"
    normalize_text(sub, font_size=15, color=RGBColor(191, 203, 222), align=PP_ALIGN.CENTER)


def optimize() -> None:
    prs = Presentation(str(SRC))
    total = len(prs.slides)
    slide_width = prs.slide_width

    for idx, slide in enumerate(prs.slides, 1):
        if idx == 1:
            style_cover(slide)
            continue
        if idx == 2:
            style_agenda(slide)
            continue
        if idx == total:
            style_end(slide)
            continue

        set_slide_bg(slide, LIGHT_BG)
        title_shape = get_title_shape(slide)
        title_text = " ".join((title_shape.text or "").split()) if title_shape else ""
        label, accent = SECTION_STYLES.get(title_text, ("Platform", TEAL))
        add_header_band(slide, label, accent, slide_width)
        add_footer(slide, idx, total, accent)

        if title_shape:
            title_shape.left = Inches(0.75)
            title_shape.top = Inches(0.52)
            title_shape.width = Inches(4.6)
            title_shape.height = Inches(0.56)
            style_title_shape(title_shape, accent)

        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text = " ".join((shape.text or "").split()) if getattr(shape, "has_text_frame", False) else ""
                if text:
                    if "圆角矩形" in shape.name or shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                        style_process_box(shape, accent)
                    elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE and len(text) <= 120:
                        style_body_text(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text = " ".join((shape.text or "").split()) if getattr(shape, "has_text_frame", False) else ""
                if text:
                    style_body_text(shape)

        # add subtle section title line
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.15), slide_width - Inches(1.5), Pt(1.5))
        line.fill.solid()
        line.fill.fore_color.rgb = LINE
        line.line.fill.background()

    prs.save(str(OUT))


if __name__ == "__main__":
    optimize()
